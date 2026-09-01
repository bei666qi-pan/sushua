from __future__ import annotations

import io
import tempfile
import unittest
from pathlib import Path
from typing import Protocol
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from document_worker.parsers import MarkItDownParser


class SavableDocument(Protocol):
    def save(self, path: Path) -> None: ...


class MarkItDownParserTests(unittest.TestCase):
    def setUp(self) -> None:
        self.parser = MarkItDownParser()

    def test_html_keeps_learning_structure_without_active_content(self) -> None:
        parsed = self.parser.parse(
            b"""
            <html>
              <head><style>.answer { color: red; }</style></head>
              <body>
                <h1>Cell biology</h1>
                <p>The mitochondrion supplies energy.</p>
                <ul><li>ATP</li><li>Respiration</li></ul>
                <script>window.stolen = document.cookie</script>
              </body>
            </html>
            """,
            "text/html",
            "html-source-sha",
        )

        self.assertEqual(parsed.parser, "markitdown")
        self.assertEqual(parsed.parser_version, "0.1.7")
        self.assertEqual(len(parsed.pages), 1)
        block = parsed.pages[0]["blocks"][0]
        self.assertEqual(block["blockType"], "text")
        self.assertIn("# Cell biology", block["markdown"])
        self.assertIn("The mitochondrion supplies energy.", block["text"])
        self.assertIn("ATP", block["text"])
        self.assertNotIn("window.stolen", block["text"])
        self.assertNotIn("color: red", block["text"])
        self.assertEqual(block["bbox"], [0, 0, 1, 1])
        self.assertLess(block["confidence"], 1)

    def test_docx_extracts_headings_paragraphs_and_tables(self) -> None:
        document = Document()
        document.add_heading("Chapter one", level=1)
        document.add_paragraph("Photosynthesis converts light energy.")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Term"
        table.cell(0, 1).text = "Meaning"
        table.cell(1, 0).text = "ATP"
        table.cell(1, 1).text = "Energy carrier"

        parsed = self.parser.parse(
            self._save_document(document, ".docx"),
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "docx-source-sha",
        )

        markdown = parsed.pages[0]["blocks"][0]["markdown"]
        self.assertIn("Chapter one", markdown)
        self.assertIn("Photosynthesis converts light energy.", markdown)
        self.assertIn("ATP", markdown)
        self.assertIn("Energy carrier", markdown)

    def test_pptx_extracts_slide_text(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Immune response"
        slide.placeholders[1].text = "Antibodies bind antigens"

        parsed = self.parser.parse(
            self._save_document(presentation, ".pptx"),
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "pptx-source-sha",
        )

        text = parsed.pages[0]["blocks"][0]["text"]
        self.assertIn("Immune response", text)
        self.assertIn("Antibodies bind antigens", text)

    def test_xlsx_extracts_sheet_names_and_cell_values(self) -> None:
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Vocabulary"
        sheet.append(["Term", "Meaning"])
        sheet.append(["Mitosis", "Cell division"])

        parsed = self.parser.parse(
            self._save_document(workbook, ".xlsx"),
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "xlsx-source-sha",
        )

        markdown = parsed.pages[0]["blocks"][0]["markdown"]
        self.assertIn("Vocabulary", markdown)
        self.assertIn("Mitosis", markdown)
        self.assertIn("Cell division", markdown)

    def test_corrupt_office_file_is_rejected(self) -> None:
        with self.assertRaisesRegex(ValueError, "invalid_document_content"):
            self.parser.parse(
                b"not an office archive",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "corrupt-source-sha",
            )

    def test_high_expansion_office_archive_is_rejected_before_conversion(self) -> None:
        output = io.BytesIO()
        with ZipFile(output, "w", compression=ZIP_DEFLATED) as archive:
            archive.writestr("word/document.xml", b"0" * (2 * 1024 * 1024))

        with self.assertRaisesRegex(ValueError, "document_archive_budget_exceeded"):
            self.parser.parse(
                output.getvalue(),
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                "expanding-source-sha",
            )

    @staticmethod
    def _save_document(document: SavableDocument, suffix: str) -> bytes:
        with tempfile.TemporaryDirectory() as directory:
            path = Path(directory) / f"source{suffix}"
            document.save(path)
            return path.read_bytes()


if __name__ == "__main__":
    unittest.main()
